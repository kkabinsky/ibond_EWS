# -*- coding: utf-8 -*-
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = pptx.Presentation("presentation.pptx")
print("Current slide count:", len(prs.slides))

blank_slide_layout = prs.slide_layouts[0]

# Theme Colors
NAVY = RGBColor(0x16, 0x27, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BLUE = RGBColor(0x0F, 0x17, 0x2A)
TEXT_MUTED = RGBColor(0x47, 0x55, 0x69)
BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER_COLOR = RGBColor(0xBF, 0xDB, 0xFE)

def add_header(slide, title_text, category_text="APPENDIX RESULTS & DATABASE SCHEMA"):
    header_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.13), Inches(0.9))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = NAVY

def style_table(table, df_data, headers):
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            
    for i, row in enumerate(df_data):
        bg = BG_LIGHT if i % 2 == 0 else BG_WHITE
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK_BLUE
                if j == 0 or "WINNER" in str(val) or "Approach" in str(val) or "109.1m" in str(val):
                    p.font.bold = True

# --- SLIDE 59: iBond 33 Features Model Comparison ---
s59 = prs.slides.add_slide(blank_slide_layout)
add_header(s59, "Table A.1: 3-Engine EWS Model Comparison (19F Base vs 33F App 1 vs 33F App 2 XGBoost)")
t_shape59 = s59.shapes.add_table(8, 5, Inches(0.6), Inches(1.5), Inches(12.13), Inches(4.8))
headers_59 = ["Metric / Evaluation Dimension", "19F Base (bond_ews.py)", "33F Approach 1 (Cox Hazard)", "33F Approach 2 (XGBoost)", "Winner & Strategic Advantage"]
data_59 = [
    ["ROC-AUC (In-Sample)", "0.9988", "0.9988", "0.9998", "Approach 2 XGBoost (Near-perfect fit)"],
    ["ROC-AUC (Out-of-Sample LOO)", "0.8640", "0.8086", "0.8871", "Approach 2 XGBoost [OVERALL WINNER 🏆]"],
    ["Recall / Sensitivity (OOS)", "100.0% (8/8)", "100.0% (8/8)", "90.42%", "Approach 1 & Base detect 100% of defaults"],
    ["Alert Volume / False Alarm Rate", "1.7% (283 rows)", "2.4% (396 rows)", "0.6% (103 rows)", "Approach 2 XGBoost [WINNER 🏆] (Lowest false alarms)"],
    ["Median Lead Time (Warning)", "326 Days (10.7m)", "404 Days (13.3m)", "152 Days (5.0m)", "Approach 2 XGBoost [WINNER 🏆] (Closest to 3M target)"],
    ["F1-Score (Out-of-Sample)", "0.5420", "0.5527", "0.6943", "Approach 2 XGBoost [WINNER 🏆] (Highest balance)"],
    ["Ease of Implementation", "High", "Very High (Linear SQL)", "High (XGBoost + SHAP)", "Approach 1 (Easiest linear SQL deployment)"],
]
style_table(t_shape59.table, data_59, headers_59)

# --- SLIDE 60: 8-Model Walk-Forward Out-of-Sample Benchmark ---
s60 = prs.slides.add_slide(blank_slide_layout)
add_header(s60, "Table A.2: 8-Model Walk-Forward Out-of-Sample Benchmark Suite")
t_shape60 = s60.shapes.add_table(9, 9, Inches(0.6), Inches(1.5), Inches(12.13), Inches(4.9))
headers_60 = ["Rank", "Model Name", "Group", "ROC-AUC", "PR-AUC", "Brier", "MCC", "Recall", "Lead Time (Median)"]
data_60 = [
    ["1", "A1-XGBoost [WINNER]", "Approach 1", "0.8314", "0.0129", "0.0050", "0.0655", "47.62%", "335 Days (11.0m)"],
    ["2", "A1-MLP (Deep Learning)", "Approach 1 (DL)", "0.8120", "0.0178", "0.1370", "0.0675", "48.81%", "610 Days (20.0m)"],
    ["3", "A2-MLP (Deep Learning)", "Approach 2 (DL)", "0.8023", "0.0150", "0.0122", "0.0800", "55.95%", "396 Days (13.0m)"],
    ["4", "A2-GRU (Sequential DL)", "Approach 2 (DL)", "0.7961", "0.0162", "0.0352", "0.0675", "48.81%", "365 Days (12.0m)"],
    ["5", "A2-Logistic", "Approach 2", "0.7707", "0.0065", "0.0073", "0.0344", "29.76%", "181 Days (6.0m)"],
    ["6", "A2-XGBoost", "Approach 2", "0.7405", "0.0102", "0.0027", "0.0717", "51.19%", "183 Days (6.0m)"],
    ["7", "A2-RandomForest", "Approach 2", "0.7223", "0.0145", "0.0031", "0.0758", "53.57%", "516 Days (17.0m)"],
    ["8", "A1-Logistic", "Approach 1", "0.6544", "0.0039", "0.8111", "0.0075", "14.29%", "534 Days (17.5m)"],
]
style_table(t_shape60.table, data_60, headers_60)

# --- SLIDE 61: Economic Loss Function Evaluation ---
s61 = prs.slides.add_slide(blank_slide_layout)
add_header(s61, "Table A.3: Economic Loss Mitigation & Financial Benefit Performance")
t_shape61 = s61.shapes.add_table(9, 7, Inches(0.6), Inches(1.5), Inches(12.13), Inches(4.9))
headers_61 = ["Model Name", "Miss Rate", "False Alarm Rate", "Protected Value (MTHB)", "Review Cost (MTHB)", "Net Benefit (MTHB)", "ROI Ratio"]
data_61 = [
    ["A1-Logistic [WINNER]", "85.71%", "9.99%", "3,510.00", "153.90", "+3,356.10 MTHB", "21.81 x"],
    ["A2-RandomForest", "46.43%", "9.88%", "2,835.00", "152.25", "+2,682.75 MTHB", "17.62 x"],
    ["A2-Logistic", "70.24%", "9.95%", "2,835.00", "153.25", "+2,681.75 MTHB", "17.50 x"],
    ["A1-XGBoost", "52.38%", "9.90%", "2,700.00", "152.50", "+2,547.50 MTHB", "16.70 x"],
    ["A1-MLP (Deep Learning)", "51.19%", "9.89%", "2,700.00", "152.45", "+2,547.55 MTHB", "16.71 x"],
    ["A2-MLP (Deep Learning)", "44.05%", "9.88%", "2,700.00", "152.15", "+2,547.85 MTHB", "16.75 x"],
    ["A2-XGBoost", "48.81%", "9.89%", "2,430.00", "152.35", "+2,277.65 MTHB", "14.95 x"],
    ["A2-GRU (Sequential DL)", "51.19%", "9.89%", "2,295.00", "152.45", "+2,142.55 MTHB", "14.05 x"],
]
style_table(t_shape61.table, data_61, headers_61)

# --- SLIDE 62: Survivor2 & DNS Factor Summary ---
s62 = prs.slides.add_slide(blank_slide_layout)
add_header(s62, "Table A.4: Survivor2 Engine & Dynamic Nelson-Siegel (DNS) Validation")
t_shape62 = s62.shapes.add_table(11, 3, Inches(0.6), Inches(1.5), Inches(12.13), Inches(4.9))
headers_62 = ["Evaluation Parameter / Metric", "Empirical Calculated Value", "Economic & Risk Operational Interpretation"]
data_62 = [
    ["Monitored Issuer Universe (N)", "289 Corporate Issuers (16,686 firm-months)", "Comprehensive coverage of all Thai corporate bond issuers"],
    ["Observed Default Events (E)", "8 Defaulted Issuers (32 positive months / 0.19%)", "Real credit default events from ThaiBMA"],
    ["Successful Early Warnings Detected", "8 / 8 Defaulted Issuers (100% Success)", "Zero missed defaulting issuers across 16,686 panel months"],
    ["Approach 2 Out-of-Sample ROC-AUC", "0.8871 (0.887)", "Highest discrimination performance on out-of-sample LOO-CV"],
    ["Approach 2 Controlled Alert Rate", "103 Firm-Months (0.6% Controlled Alert Rate)", "Ultra-precise risk classification (lowest false alarm volume)"],
    ["Approach 2 Median Lead Time", "152.0 Days (~5.0 Months)", "Closest lead time buffer to target 3-month window"],
    ["DNS Level Factor (beta_1) Correlation", "+0.962", "Strong alignment with long-term 15-year government bond yield"],
    ["DNS Slope Factor (beta_2) Correlation", "+0.987", "Strong alignment with term structure slope (10Y - 1Y spread)"],
    ["DNS Curvature Factor (beta_3) Correlation", "+0.905", "Captures yield curve hump/twist (2*2Y - 3M - 10Y)"],
    ["DNS Overall Curve Fitting Error", "RMSE = 0.0384%", "Ultra-high precision yield curve fitting across 3M to 10Y tenors"],
]
style_table(t_shape62.table, data_62, headers_62)

# --- SLIDE 63: Top Monitored iBond Corporate Bonds Lead Time ---
s63 = prs.slides.add_slide(blank_slide_layout)
add_header(s63, "Table A.5: Top Distressed iBond Corporate Bonds Lead Time Breakdown")
t_shape63 = s63.shapes.add_table(11, 9, Inches(0.6), Inches(1.5), Inches(12.13), Inches(4.9))
headers_63 = ["Symbol", "Issuer", "Company Name", "Sector", "19F Base", "33F App 1", "33F App 2 XGB", "XGB PD3M", "Lead Time Winner"]
data_63 = [
    ["ECF266A", "ECF", "East Coast Furnitech", "INDUS", "120d (3.9m)", "158d (5.2m)", "97d (3.2m)", "0.577", "App 2 XGBoost 🏆 (3.2m)"],
    ["SQ266A", "SQ", "Sahakol Equipment PCL", "INDUS", "370d (12.2m)", "474d (15.6m)", "109d (3.6m)", "0.543", "App 2 XGBoost 🏆 (3.6m)"],
    ["PRIME25NA", "PRIME", "Prime Road Power PCL", "ENERG", "310d (10.2m)", "391d (12.8m)", "118d (3.9m)", "0.465", "App 2 XGBoost 🏆 (3.9m)"],
    ["TPOLY266A", "TPOLY", "Thai Poly Cons. PCL", "CONS", "380d (12.5m)", "485d (15.9m)", "148d (4.9m)", "0.223", "App 2 XGBoost 🏆 (4.9m)"],
    ["A24NA", "A", "Areeya Property PCL", "PROP", "290d (9.5m)", "368d (12.1m)", "156d (5.1m)", "0.448", "App 2 XGBoost 🏆 (5.1m)"],
    ["GRAND254A", "GRAND", "Grande Asset Hotels", "PROP", "330d (10.9m)", "418d (13.7m)", "206d (6.8m)", "0.022", "App 2 XGBoost 🏆 (6.8m)"],
    ["JCK266A", "JCK", "JCK International", "PROP", "320d (10.5m)", "416d (13.7m)", "232d (7.6m)", "0.082", "App 2 XGBoost 🏆 (7.6m)"],
    ["PF265A", "PF", "Property Perfect PCL", "PROP", "290d (9.5m)", "370d (12.2m)", "278d (9.1m)", "0.337", "App 2 XGBoost 🏆 (9.1m)"],
    ["STARK242A", "STARK", "Stark Corporation", "INDUS", "280d (9.2m)", "364d (12.0m)", "180d (5.9m)", "0.412", "App 2 XGBoost 🏆 (5.9m)"],
    ["ALL241A", "ALL", "All Inspire Dev.", "PROP", "250d (8.2m)", "320d (10.5m)", "160d (5.2m)", "0.380", "App 2 XGBoost 🏆 (5.2m)"],
]
style_table(t_shape63.table, data_63, headers_63)

# --- SLIDE 64: Head-to-Head Performance Comparison Across 3 Engines ---
s64 = prs.slides.add_slide(blank_slide_layout)
add_header(s64, "Table A.6: Ultimate 3-Engine EWS Performance Matrix (19F vs App 1 vs App 2 XGBoost)")
t_shape64 = s64.shapes.add_table(4, 8, Inches(0.6), Inches(1.5), Inches(12.13), Inches(4.5))
headers_64 = ["EWS Engine & Model Type", "Features", "AUC (In-Sample)", "AUC (Out-of-Sample LOO)", "Recall (%)", "Alert Rate (%)", "Median Lead Time", "Operational Verdict"]
data_64 = [
    ["Baseline bond_ews.py", "19 Features", "0.9988", "0.8640", "100.0% (8/8)", "1.7% (283 rows)", "326 Days (10.7m)", "Baseline 19-Feature Model"],
    ["Approach 1: Cox Hazard (run_survivor_ews_33features.py)", "33 Features", "0.9988", "0.8086", "100.0% (8/8)", "2.4% (396 rows)", "404 Days (13.3m)", "Long-term Linear Decay Model"],
    ["Approach 2: Calibrated XGBoost (run_survivor_ews_33features_xgb.py)", "33 Features", "0.9998", "0.8871", "90.42%", "0.6% (103 rows)", "152 Days (5.0m)", "OVERALL WINNER 🏆 (Ultra Precise)"],
]
style_table(t_shape64.table, data_64, headers_64)

prs.save("presentation_final.pptx")
print("Saved updated presentation to presentation_final.pptx. Total slides:", len(prs.slides))
try:
    import shutil
    shutil.copyfile("presentation_final.pptx", "presentation.pptx")
    print("Successfully overwritten presentation.pptx")
except Exception:
    print("presentation_final.pptx has been saved with the latest slides.")
